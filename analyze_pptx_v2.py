from pptx import Presentation
import sys

prs = Presentation(r'C:\Users\Avalon\Desktop\Manhattan_Protocol_MEL_Brief.pptx')

print("\n" + "="*100)
print("MANHATTAN PROTOCOL PRESENTATION ANALYSIS")
print("="*100)

print(f"\nBasic Info:")
print(f"  Total Slides: {len(prs.slides)}")
print(f"  Slide Dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
print(f"  Aspect Ratio: 16:9")

print(f"\n" + "="*100)
print("SLIDE 1 - COVER SLIDE STRUCTURE")
print("="*100)

slide1 = prs.slides[0]
print(f"Total Shapes: {len(slide1.shapes)}\n")

for i, shape in enumerate(slide1.shapes[:15]):
    if hasattr(shape, "text_frame") and shape.text:
        text = shape.text.replace("\n", " // ")[:80]
        print(f"  {i:2d}. {text}")
        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                size = r.font.size.pt if r.font.size else "default"
                bold = "BOLD" if r.font.bold else ""
                font = r.font.name or "default"
                print(f"       Font: {font} {size}pt {bold}")

print(f"\n" + "="*100)
print("SLIDE 2 - CONTENT SLIDE STRUCTURE")
print("="*100)

slide2 = prs.slides[1]
print(f"Total Shapes: {len(slide2.shapes)}\n")

sections = []
for i, shape in enumerate(slide2.shapes):
    if hasattr(shape, "text_frame") and shape.text:
        text = shape.text.strip()
        if text and len(text) < 100:
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    r = p.runs[0]
                    size = r.font.size.pt if r.font.size else "default"
                    bold = r.font.bold
                    if bold or size >= 13:
                        print(f"  [HEADER] {text}")
                        print(f"           Font: {r.font.name} {size}pt")
        elif text and len(text) < 200:
            print(f"  [BODY] {text[:120]}")

print(f"\n" + "="*100)
print("STYLE OBSERVATIONS")
print("="*100)

fonts_used = set()
sizes_used = set()

for slide in prs.slides[:5]:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.name:
                        fonts_used.add(run.font.name)
                    if run.font.size:
                        sizes_used.add(run.font.size.pt)

print(f"\nFonts Used: {sorted(list(fonts_used))}")
print(f"Font Sizes: {sorted(list(sizes_used))}")

print(f"\n" + "="*100)
print("CONTENT PATTERN")
print("="*100)

for slide_num, slide in enumerate(prs.slides[:3], 1):
    text_shapes = [s for s in slide.shapes if hasattr(s, "text_frame") and s.text.strip()]
    print(f"\nSlide {slide_num}: {len(text_shapes)} text elements")

    for shape in text_shapes[:8]:
        text = shape.text.strip()[:60]
        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                bold = "**" if r.font.bold else ""
                print(f"  {bold}{text}{bold}")
