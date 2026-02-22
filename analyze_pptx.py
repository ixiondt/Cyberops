from pptx import Presentation
import json

prs = Presentation(r'C:\Users\Avalon\Desktop\Manhattan_Protocol_MEL_Brief.pptx')

analysis = {
    "total_slides": len(prs.slides),
    "slide_dimensions": f"{prs.slide_width.inches}\" x {prs.slide_height.inches}\"",
    "slides": []
}

for i, slide in enumerate(prs.slides):
    slide_info = {
        "slide_number": i + 1,
        "shapes": len(slide.shapes),
        "elements": []
    }

    for shape in slide.shapes:
        element = {
            "type": shape.shape_type,
            "has_text": hasattr(shape, "text_frame")
        }

        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            element["text_length"] = len(shape.text)
            element["text_preview"] = shape.text[:100] if shape.text else ""
            element["paragraphs"] = len(text_frame.paragraphs)

            # Font info from first paragraph
            if text_frame.paragraphs:
                p = text_frame.paragraphs[0]
                if p.runs:
                    run = p.runs[0]
                    element["font_name"] = run.font.name
                    element["font_size"] = run.font.size.pt if run.font.size else "default"
                    element["bold"] = run.font.bold
                    element["color"] = str(run.font.color.rgb) if run.font.color.type else "none"
                element["alignment"] = str(p.alignment)

        if hasattr(shape, "fill"):
            element["fill_type"] = str(shape.fill.type)

        slide_info["elements"].append(element)

    analysis["slides"].append(slide_info)

# Pretty print
for slide in analysis["slides"][:5]:
    print(f"\n{'='*80}")
    print(f"SLIDE {slide['slide_number']} - {slide['shapes']} shapes")
    print(f"{'='*80}")
    for i, elem in enumerate(slide['elements']):
        print(f"\nElement {i+1}:")
        print(f"  Type: {elem.get('type')}")
        print(f"  Has Text: {elem.get('has_text')}")
        if elem.get('text_preview'):
            print(f"  Text: {elem['text_preview']}")
        if elem.get('font_name'):
            print(f"  Font: {elem['font_name']} {elem['font_size']}pt (bold={elem['bold']})")
        print(f"  Fill: {elem.get('fill_type')}")

print(f"\n\n{'='*80}")
print(f"SUMMARY")
print(f"{'='*80}")
print(f"Total Slides: {analysis['total_slides']}")
print(f"Slide Dimensions: {analysis['slide_dimensions']}")
