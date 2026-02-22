#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Army Green color scheme
SAGE_GREEN = RGBColor(108, 135, 93)
DARK_GREEN = RGBColor(51, 76, 51)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(prs, title, subtitle, date_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SAGE_GREEN

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_frame.text = date_text if date_text else "UNCLASSIFIED // FOR OFFICIAL USE ONLY"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = WHITE
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = SAGE_GREEN
    header_shape.line.color.rgb = DARK_GREEN

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GREEN
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(9.4), Inches(0.2))
    footer_frame = footer_box.text_frame
    footer_frame.text = "UNCLASSIFIED // FOR OFFICIAL USE ONLY"
    footer_frame.paragraphs[0].font.size = Pt(8)
    footer_frame.paragraphs[0].font.color.rgb = DARK_GREEN

def add_section_slide(prs, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GREEN

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = section_title
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Build slides
add_title_slide(prs, "MANHATTAN PROTOCOL\nOPERATION BRIEF",
    "CPT 173 Deployment & Integration Phase\nDefensive Cyberspace Operations",
    "DTG: 171001Z FEB 2026")

add_content_slide(prs, "EXECUTIVE SUMMARY (BLUF)", [
    "Operation: Manhattan Protocol (OPORD 26-02) - DCO-M at BPEnergy",
    "Mission: Detect, contain, mitigate APT41 activity targeting CI networks",
    "Threat: APT41 (state-sponsored, PRC) - credential theft, lateral movement",
    "CPT 173 Role: Lead element - MOC operations, threat hunting, IR",
    "Phase I: Deployment & Integration (17-22 FEB) - Stand up MOC, baseline",
    "Status: READY TO EXECUTE | Risk Level: MODERATE"])

add_section_slide(prs, "AREA OF OPERATIONS")

add_content_slide(prs, "Area of Operations: BPEnergy", [
    "Primary AO: BPEnergy Enterprise Network (SI-EN)",
    "  - Microsoft AD/Exchange, corporate applications, user systems",
    "  - Headquarters: Walton, NY | Supporting 125,000 citizens",
    "Secondary AO: Cloud Environments (Azure Gov, AWS GovCloud)",
    "  - R&D data, sensitive customer information",
    "Tertiary AO: Operational Technology (OT) Enclaves",
    "  - Critical infrastructure support (power, water, etc.)",
    "  - PASSIVE MONITORING ONLY in Phase I"])

add_section_slide(prs, "THREAT ASSESSMENT")

add_content_slide(prs, "Threat: APT41 (State-Sponsored, PRC)", [
    "Classification: Tier-1 Advanced Persistent Threat",
    "  - Well-resourced with state backing | Highly adaptive",
    "Operational Model: Dual-mission (espionage + cybercrime)",
    "  - Long-term strategic access objectives",
    "Capabilities: Zero-day exploits | Supply-chain compromise",
    "  - Cloud IAM exploitation | Advanced lateral movement",
    "Infrastructure: PRC-based + global proxy network"])

add_content_slide(prs, "APT41: Attack Chain & Malware", [
    "Common Malware: ShadowPad, Winnti, Crosswalk, Cobalt Strike",
    "Attack Chain:",
    "  - Initial Access: Spearphishing, vulnerability, supply-chain",
    "  - Persistence: Scheduled tasks, registry mods, cloud tokens",
    "  - Lateral Movement: SMB, RDP, WinRM, SSH, Kerberos",
    "  - Collection: Targeted R&D data via HTTPS/DNS tunneling",
    "Defense Evasion: Log tampering, LOLBins, encrypted C2"])

add_content_slide(prs, "APT41 Threat COAs", [
    "MLCOA (Most Likely COA):",
    "  - Credential harvesting, lateral movement, data exfiltration",
    "  - Establish dormant persistence for long-term access",
    "MDCOA (Most Dangerous COA):",
    "  - Compromise OT networks supporting defense production",
    "  - Deploy destructive malware to disrupt manufacturing",
    "  - Simultaneous cloud + on-premises attacks"])

add_section_slide(prs, "MISSION & INTENT")

add_content_slide(prs, "Mission Statement", [
    "Deploy to Bakerville NLT 23 FEB 2026 to conduct Defensive",
    "Cyberspace Operations to identify, contain, mitigate malicious",
    "cyber activity targeting BPEnergy's CI networks.",
    "",
    "End State:",
    "BPEnergy networks hardened, adversary presence removed or",
    "contained, supported organization sustains secure operations."])

add_content_slide(prs, "Commander's Intent (CPT 173 Focus)", [
    "Purpose: Protect BPEnergy critical networks from adversary activity",
    "Key Tasks:",
    "  - Establish primary Mission Operations Center (MOC)",
    "  - Conduct enterprise-wide vulnerability assessment",
    "  - Lead incident response operations",
    "  - Maintain continuous monitoring and daily reporting",
    "Constraints:",
    "  - No production system changes without explicit approval",
    "  - No direct OT actions without BPEnergy coordination"])

add_section_slide(prs, "PHASE I: DEPLOYMENT")

add_content_slide(prs, "Phase I Objectives (17-22 FEB 2026)", [
    "1. Establish secure MOC with 24/7 operational capability",
    "2. Validate network access (SIEM, EDR, cloud platforms)",
    "3. Confirm authorities and ROE (in writing)",
    "4. Complete baseline network assessment",
    "5. Develop and validate threat hunting hypotheses",
    "6. Deploy detection rules and playbooks",
    "7. Integrate CPT 173 with BPEnergy cybersecurity staff"])

add_content_slide(prs, "Phase I Timeline", [
    "Day 1-2 (17-18 FEB): ARRIVAL & I-TEM",
    "  - Personnel deployment, secure comms, I-TEM with leadership",
    "Day 3-4 (19-20 FEB): MOC SETUP & BASELINE",
    "  - 24/7 MOC operations, log analysis, vulnerability scan",
    "Day 5-6 (21-22 FEB): HARDENING & PHASE II PREP",
    "  - Final assessments, Initial Assessment Report",
    "  - ARCYBER approval for Phase II transition"])

add_section_slide(prs, "MOC OPERATIONS")

add_content_slide(prs, "MOC Organization (24/7)", [
    "CPT 173 Task Organization:",
    "  - Team Lead (CPT) - Overall leadership, ARCYBER coordination",
    "  - Deputy TL (WO1/CW2) - Backup, forensics oversight",
    "  - Forensics Lead (SSG/SFC) - Host analysis, artifacts",
    "  - Network Analysis Lead (SSG/SFC) - Traffic analysis, hunting",
    "  - Incident Response Lead (SSG/SFC) - Playbooks, containment",
    "  - Reporting NCO (SGT) - SITREPs, hunt logs"])

add_content_slide(prs, "MOC Functions", [
    "Synchronize CPT 173 activities (daily syncs 0800Z/2000Z)",
    "Maintain Common Operational Picture (COP)",
    "Manage reporting to ARCYBER and BPEnergy",
    "Coordinate DCO-RA requests and escalation",
    "Maintain mission logs and battle rhythm",
    "Provide 24/7 incident detection and triage"])

add_section_slide(prs, "THREAT HUNTING")

add_content_slide(prs, "Threat Hunting Hypotheses (Phase II)", [
    "Hypothesis 1: APT41 targets R&D cloud via credential spray",
    "  - Look for: Multiple failed logins, unusual IPs, MFA anomalies",
    "Hypothesis 2: APT41 uses supply-chain vendor access",
    "  - Look for: Service account anomalies, vendor segments accessed",
    "Hypothesis 3: APT41 attempts OT reconnaissance",
    "  - Look for: Enterprise hosts querying OT, PLC discovery",
    "Hypothesis 4: APT41 establishes persistence via tasks/services",
    "  - Look for: New scheduled tasks, services, registry mods"])

add_content_slide(prs, "Detection Rules Deployed", [
    "Network-Level:",
    "  - C2 beaconing patterns, data exfiltration, lateral movement",
    "Host-Level:",
    "  - Process anomalies, credential access, persistence mechanisms",
    "Cloud-Level:",
    "  - IAM anomalies, suspicious API calls, storage enumeration",
    "All rules tested and ready for Phase II active hunting"])

add_section_slide(prs, "PIRs & REQUIREMENTS")

add_content_slide(prs, "Priority Intelligence Requirements (PIRs)", [
    "PIR 1: Has APT41 established persistence in BPEnergy?",
    "  Status: UNKNOWN - no indicators found yet",
    "PIR 2: What TTPs is APT41 currently employing?",
    "  Status: PARTIAL - historical TTPs known",
    "PIR 3: Is APT41 accessing OT systems?",
    "  Status: UNKNOWN - OT baseline established",
    "PIR 4: What indicators reveal imminent exfiltration?",
    "  Status: DETECTION RULES DEPLOYED"])

add_section_slide(prs, "RESOURCES")

add_content_slide(prs, "Critical Resources & Authorities", [
    "Personnel: CPT 173 (6) + Battalion HQ support (3) on-site",
    "Data Access:",
    "  - SIEM read-only (past 90 days) | EDR agents authorized",
    "  - Azure/AWS read-only (audit logs) | DNS/proxy logs",
    "  - OT network monitoring (passive only)",
    "Authorities:",
    "  - Title 10, DCO-M via ARCYBER FRAGO 01-26 - CONFIRMED",
    "  - DCO-RA (explicit approval required) - CONFIRMED"])

add_content_slide(prs, "Phase I Key Decision Points", [
    "DP-1: Communications operational - 17 FEB 1200Z",
    "DP-2: Network architecture understood - 18 FEB 1100Z",
    "DP-3: Access procedures approved - 18 FEB 1200Z",
    "DP-4: Authorities confirmed in writing - 18 FEB 1300Z",
    "DP-5: All access validated & operational - 18 FEB 1700Z",
    "DP-6: Playbooks reviewed & trained - 20 FEB 1900Z",
    "DP-7: ARCYBER approves Phase II - 21 FEB 1600Z"])

add_section_slide(prs, "RISK ASSESSMENT")

add_content_slide(prs, "Identified Risks & Mitigations", [
    "Risk 1: Access Delays - PROBABILITY: MODERATE | IMPACT: MODERATE",
    "  Mitigation: Establish early; escalate via ARCYBER",
    "Risk 2: APT41 Active in Phase I - PROBABILITY: MODERATE | IMPACT: HIGH",
    "  Mitigation: Pre-position playbooks; 24/7 MOC; rapid escalation",
    "Risk 3: OT Disruption - PROBABILITY: LOW | IMPACT: CRITICAL",
    "  Mitigation: NO direct OT actions; passive monitoring only",
    "Risk 4: Intelligence Gaps - PROBABILITY: HIGH | IMPACT: MODERATE",
    "  Mitigation: Leverage 91st CPB reachback; aggressive hunting"])

add_section_slide(prs, "QUESTIONS?")

# Save presentation
output_path = "C:\\Users\\Avalon\\Desktop\\CPT173_PHASE_I_BRIEF.pptx"
prs.save(output_path)
print(f"SUCCESS: PowerPoint created at {output_path}")
print(f"Total slides: {len(prs.slides)}")
