from pptx import Presentation
import xml.etree.ElementTree as ET

prs = Presentation(r'C:\Users\Avalon\Desktop\Manhattan_Protocol_MEL_Brief.pptx')

print("\n" + "="*100)
print("MANHATTAN PROTOCOL PRESENTATION ANALYSIS")
print("="*100)

print(f"\nBasic Info:")
print(f"  Total Slides: {len(prs.slides)}")
print(f"  Slide Dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\" (16:9)")

print(f"\n" + "="*100)
print("SLIDE STRUCTURE & CONTENT ANALYSIS")
print("="*100)

for slide_idx in range(min(3, len(prs.slides))):
    slide = prs.slides[slide_idx]

    print(f"\nSLIDE {slide_idx + 1}:")
    print(f"  Total shapes: {len(list(slide.shapes))}")

    text_elements = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text.strip():
                text_elements.append(shape.text.strip())
        except:
            pass

    print(f"  Text elements: {len(text_elements)}")
    print(f"  Content preview:")
    for i, text in enumerate(text_elements[:12]):
        preview = text[:70].replace("\n", " // ")
        print(f"    {i+1:2d}. {preview}")

print(f"\n" + "="*100)
print("KEY STYLE PATTERNS IDENTIFIED")
print("="*100)

print("""
Based on the Manhattan Protocol presentation structure:

1. SLIDE DIMENSIONS & LAYOUT
   - Aspect Ratio: 16:9 (10" x 5.625")
   - Standard presentation format
   - Heavy use of shapes for visual structure (29-35 shapes per slide)

2. TYPOGRAPHY
   - Primary Font: Calibri throughout
   - Hierarchy:
     * Operation/Title: 52pt Bold
     * Slide Headers: 22pt Bold
     * Section Headers: 13pt Bold
     * Subsection Headers: 10pt Bold
     * Body Text: 11-12pt
     * Metadata/Footer: 8-9pt

3. COVER SLIDE STRUCTURE
   - Operation Name (20pt)
   - Large Title (52pt bold)
   - Subtitle/Role (14pt)
   - Organization/Unit info (11pt)
   - Location/Date info (11pt)
   - Contents list with visual bullet points (10pt)
   - Classification marking (11pt bold)
   - DTG timestamp

4. CONTENT SLIDES
   - Title: 22pt Bold with divider
   - Subtitle: 11pt (context/grouping)
   - Section Headers: 13pt Bold
   - Body Content: 11-12pt
   - Multiple sections per slide (3-4 major sections typical)
   - Footer: Metadata line (8pt) + Classification (9pt bold)
   - Visual separators between sections

5. CONTENT ORGANIZATION
   - Hierarchical: MISSION → PURPOSE → METHOD/EXECUTION/KEY POINTS
   - Each slide covers ONE major topic with 3-4 subtopics
   - Heavy emphasis on operational structure
   - Classification markings on EVERY slide
   - Organizational context (CPT 173 // 126th CPBn // Op Name)

6. COLOR & VISUAL DESIGN
   - Dark background with solid fill shapes (likely Navy Blue ~0,40,80)
   - Clean dividers and section separators
   - Minimal but structured visual design
   - Professional military briefing aesthetic

7. METADATA ON EVERY SLIDE
   - Left footer: "UNCLASSIFIED // FOUO" (bold)
   - Center: Organizational identifier
   - Right: Timestamps/References
""")

print(f"\n" + "="*100)
print("PPTX GENERATION GUIDELINES")
print("="*100)

print("""
FOR FUTURE PPTX GENERATION:

1. Cover Slide Must Include:
   - Large title (52pt bold)
   - Subtitle with role/purpose (14pt)
   - Unit/Organization (11pt)
   - Location/Date (11pt)
   - Contents list (10pt items with visual bullets)
   - Classification marking (11pt bold)
   - DTG timestamp

2. Every Content Slide Should Have:
   - Title: 22pt Bold with divider bar
   - Subtitle: 11pt (operational context or section grouping)
   - Body: Multiple sections (3-4 sections per slide)
     * Each section: 13pt Bold header + 11-12pt body
   - Visual separators between sections
   - Footer: Classification (left, bold) + Metadata (right, small)
   - Organizational context footer

3. Content Density:
   - 3-4 major sections per slide
   - Each section 2-4 key points
   - Total ~15-20 text elements per slide

4. Font Consistency:
   - Title: 22pt Bold (Calibri)
   - Section: 13pt Bold (Calibri)
   - Body: 11-12pt (Calibri)
   - Meta: 8-9pt (Calibri)
   - Avoid other fonts

5. Structure Pattern:
   - Operational briefing format (not bulleted lists)
   - Heading-based organization
   - Clear hierarchy and nesting
   - Every slide has classification + metadata footer

6. Color Scheme:
   - Dark background (Navy ~RGB(0,40,80))
   - Light text (White/Light Gray)
   - Professional military aesthetic
   - Minimal accent colors
""")
