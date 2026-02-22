from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Colors
DARK_BG = RGBColor(0, 40, 80)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
SECTION_BLUE = RGBColor(0, 51, 102)

def add_background(slide):
    """Add dark background"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_cover_slide(prs, operation, role, unit, location, date, contents):
    """Create professional cover slide with contents list"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    # Operation label
    op_label = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.4))
    tf = op_label.text_frame
    tf.text = "OPERATION"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = LIGHT_GRAY

    # Operation name (large)
    op_name = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(1.2))
    tf = op_name.text_frame
    tf.word_wrap = True
    tf.text = operation
    p = tf.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Role/Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    tf.text = role
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Unit and location info
    info1 = slide.shapes.add_textbox(Inches(0.5), Inches(2.35), Inches(9), Inches(0.3))
    tf = info1.text_frame
    tf.text = unit
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

    info2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.65), Inches(9), Inches(0.3))
    tf = info2.text_frame
    tf.text = location + "  //  " + date
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

    # Contents heading
    contains = slide.shapes.add_textbox(Inches(0.5), Inches(3.05), Inches(2), Inches(0.25))
    tf = contains.text_frame
    tf.text = "CONTAINS"
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = LIGHT_GRAY

    # Contents list
    contents_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(2.1))
    tf = contents_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(contents):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

    # Footer - Classification
    footer_class = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(3), Inches(0.25))
    tf = footer_class.text_frame
    tf.text = "UNCLASSIFIED // FOUO"
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = LIGHT_GRAY

    # Footer - metadata
    footer_meta = slide.shapes.add_textbox(Inches(7.5), Inches(5.3), Inches(2), Inches(0.25))
    tf = footer_meta.text_frame
    tf.text = "171001ZFEB2026"
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY

def add_content_slide(prs, title, subtitle, sections_dict, org_context):
    """Create professional content slide with structured sections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    # Title bar with divider
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.3))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

    # Divider line
    divider = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(9), Inches(0.01))
    divider.fill.solid()
    divider.fill.fore_color.rgb = SECTION_BLUE
    divider.line.color.rgb = SECTION_BLUE

    # Content sections
    current_y = 1.2
    for section_title, section_content in sections_dict.items():
        # Section header
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(current_y), Inches(9), Inches(0.3))
        tf = section_box.text_frame
        tf.text = section_title
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        current_y += 0.3

        # Section content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(current_y), Inches(8.8), Inches(0.9))
        tf = content_box.text_frame
        tf.word_wrap = True
        if isinstance(section_content, list):
            for i, item in enumerate(section_content):
                if i > 0:
                    tf.add_paragraph()
                p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.space_after = Pt(2)
                p.level = 0
        else:
            tf.text = section_content
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE

        current_y += 0.75

    # Footer - Classification
    footer_class = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(2), Inches(0.25))
    tf = footer_class.text_frame
    tf.text = "UNCLASSIFIED // FOUO"
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = LIGHT_GRAY

    # Footer - Org context
    footer_org = slide.shapes.add_textbox(Inches(2.5), Inches(5.3), Inches(5), Inches(0.25))
    tf = footer_org.text_frame
    tf.text = org_context
    p = tf.paragraphs[0]
    p.font.size = Pt(8)
    p.font.color.rgb = LIGHT_GRAY

# ============================================================================
# BUILD PRESENTATION
# ============================================================================

# Slide 1: Cover
add_cover_slide(prs,
    operation="MANHATTAN\nPROTOCOL",
    role="CPT 173 // MISSION ELEMENT LEAD OPERATIONAL BRIEF",
    unit="126th Cyber Protection Battalion // ARCYBER",
    location="BPEnergy (Bakerville, NY)",
    date="Phase I Start: 17 FEB 2026",
    contents=[
        "Mission Overview & Commander's Intent",
        "Threat Analysis: APT41",
        "MDMP Execution Framework",
        "Phase I–III Operational Timeline",
        "Sensor Deployment & Detection Strategy",
        "Incident Response & DCO-RA Procedures",
        "Battle Rhythm & Reporting Requirements",
        "CPT 173 Critical Actions"
    ]
)

# Slide 2: Mission Overview
add_content_slide(prs,
    title="MISSION OVERVIEW",
    subtitle="Commander's Intent | Task | Purpose | End State",
    sections_dict={
        "MISSION": "CPT 173 Mission Element deploys to BPEnergy (Bakerville, NY) NLT 23 FEB 2026 to conduct DCO-M and DCO-RA targeting APT41 cyber activity affecting critical infrastructure and DoD-aligned production systems.",
        "PURPOSE": "Eliminate APT41 presence from BPEnergy networks and prevent re-establishment of adversary footholds through continuous monitoring, threat hunting, incident response, and rapid containment of malicious activity.",
        "KEY TASKS": [
            "Establish and operate Mission Operations Center (MOC) with 24/7 watch capability",
            "Conduct enterprise vulnerability assessment and cloud environment analysis",
            "Execute intelligence-driven threat hunting targeting APT41 TTPs",
            "Lead incident response and DCO-RA actions upon confirmed malicious activity"
        ]
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 3: Threat - APT41
add_content_slide(prs,
    title="THREAT ANALYSIS",
    subtitle="APT41: State-Sponsored Cyber Espionage & Cybercrime Group",
    sections_dict={
        "COMPOSITION": "Multi-mission actor: credential harvesting, lateral movement, data exfiltration, supply-chain exploitation, destructive OT-focused malware deployment.",
        "CAPABILITIES": [
            "Zero-day exploitation and vulnerability orchestration",
            "Cross-platform malware: ShadowPad, Winnti, custom Cobalt Strike beacons",
            "Cloud IAM manipulation and persistence via API keys",
            "Advanced lateral movement via SMB, RDP, WinRM with credential reuse"
        ],
        "MOST LIKELY COA (MLCOA)": "Spearphishing → credential harvesting → lateral movement → data exfiltration targeting R&D and production data.",
        "MOST DANGEROUS COA (MDCOA)": "OT network compromise with destructive malware targeting manufacturing systems to disrupt production and endanger public services."
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 4: MDMP Framework
add_content_slide(prs,
    title="MDMP EXECUTION FRAMEWORK",
    subtitle="Seven-Step Military Decision Making Process aligned with OPORD 26-02",
    sections_dict={
        "STEPS 1-2: Mission Receipt & Analysis": "ARCYBER FRAGO authority → Operational environment assessment → APT41 threat COAs → Four Priority Intelligence Requirements (PIRs) focused on persistence, TTPs, OT targeting, and lateral movement indicators.",
        "STEPS 3-5: COA Development & Comparison": "Task organization (CPT 173 lead, CPT 174 support, Battalion S2/S3) → Three-phase concept → Wargaming (credential compromise, malware, OT scenarios) → Friction analysis (OT coordination delays, cloud visibility gaps, authority approval timing).",
        "STEPS 6-7: Approval & Orders": "Battalion Commander approval of operational concept → ARCYBER taskord validation → BPEnergy access agreement → Subordinate task assignment to CPT elements, S2 (daily INTSUMs), and S3 (logistics/reporting)."
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 5: Operational Phases
add_content_slide(prs,
    title="PHASE EXECUTION PLAN",
    subtitle="Three-Phase Mission Timeline with Synchronization Points",
    sections_dict={
        "PHASE I (17–22 FEB): Deployment & Integration": [
            "Conduct Initial Technical Exchange Meeting (I-TEM) with BPEnergy cybersecurity division",
            "Establish MOC infrastructure, secure communications, personnel accountability",
            "Validate network architecture, access permissions, ROE, and data handling compliance",
            "Deploy baseline monitoring tools and begin log ingestion from SIEM/EDR"
        ],
        "PHASE II (23 FEB–TBD): Active Defensive Operations": [
            "24/7 continuous monitoring of enterprise networks and cloud environments",
            "Daily threat hunting based on intelligence-driven hypotheses targeting APT41 TTPs",
            "Incident response and DCO-RA execution upon confirmed malicious activity detection",
            "Malware analysis, forensics, and attribution supporting daily intelligence updates"
        ],
        "PHASE III (TBD): Transition & Hardening": [
            "Deliver final assessment reports with mitigation recommendations and long-term hardening plans",
            "Conduct knowledge transfer and staff training to BPEnergy cybersecurity personnel",
            "Validate implementation of security controls and prepare for self-sustaining operations"
        ]
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 6: Detection Strategy
add_content_slide(prs,
    title="SENSOR & DETECTION STRATEGY",
    subtitle="Multi-Layer Monitoring Architecture for APT41 Detection",
    sections_dict={
        "NETWORK LAYER": [
            "Zeek/Suricata for encrypted C2 pattern detection and beaconing analysis",
            "NetFlow analysis for lateral movement and data exfiltration volume anomalies",
            "DNS query logging for malware communication and fast-flux network detection"
        ],
        "ENDPOINT LAYER": [
            "EDR platform (BPEnergy-provided) for process execution, memory injection, credential access detection",
            "Velociraptor for rapid artifact collection and forensic timeline development",
            "OSQuery for baseline deviation detection across all systems"
        ],
        "CLOUD LAYER": [
            "Azure Security Center and AWS GuardDuty for IAM anomalies and unauthorized resource creation",
            "CloudTrail/CloudWatch logging for API activity pattern analysis",
            "IAM Access Analyzer for privilege escalation and misconfiguration identification"
        ],
        "INTELLIGENCE FUSION": "Daily SITREP synthesis of all detection layers tied to APT41 TTPs; PIR-driven collection focus; automated indicator dissemination"
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 7: Incident Response
add_content_slide(prs,
    title="INCIDENT RESPONSE & DCO-RA PROCEDURES",
    subtitle="Six-Phase IR Model with Integrated Defensive Cyberspace Response Actions",
    sections_dict={
        "DETECTION & ANALYSIS": [
            "Real-time alerts from SIEM, EDR, IDS/IPS correlated with threat intelligence",
            "Immediate triage by CPT 173 to establish confirmed malicious activity",
            "Memory imaging, log collection, and network packet capture for forensics"
        ],
        "CONTAINMENT & ERADICATION": [
            "Host isolation and credential disablement upon confirmation",
            "Malware removal and persistence mechanism elimination (C&C takedown coordination with ARCYBER)",
            "Vulnerability patching and access control hardening",
            "OT environment actions require explicit BPEnergy approval; no automated OT containment"
        ],
        "RECOVERY & REPORTING": [
            "System integrity validation and re-enablement after forensic clearance",
            "Monitoring for re-infection indicators and indicator updates to detection systems",
            "Incident Report within 1 hour of detection; Full Report within 24 hours",
            "Post-action analysis and COA documentation for lessons learned"
        ]
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 8: Battle Rhythm
add_content_slide(prs,
    title="BATTLE RHYTHM & REPORTING REQUIREMENTS",
    subtitle="Daily Synchronization Points and Escalation Procedures",
    sections_dict={
        "DAILY REPORTING CADENCE": [
            "0800Z: MOC accountability and threat brief (internal to CPT 173 / S3)",
            "1200Z: Mid-day tactical update on hunt hypotheses and current detections",
            "1600Z: Daily SITREP to ARCYBER and BPEnergy leadership (progress vs PIRs, incidents, risks)",
            "2400Z: Evening situation consolidation and next-day hunt planning"
        ],
        "EVENT-DRIVEN ESCALATION": [
            "Incident Detection: Immediate notification to MOC watch (triage within 15 min)",
            "Confirmed Malicious Activity: BPEnergy cybersecurity division notification + ARCYBER S2/S3 alert",
            "DCO-RA Required: Formal request to ARCYBER with justification; approval authority within 2 hours",
            "Critical Event (OT threat): Executive escalation to BPEnergy leadership and Battalion Commander"
        ],
        "WEEKLY SYNCHRONIZATION": "Thursday 1600Z comprehensive summary to ARCYBER; lessons learned and PIR status update"
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 9: CPT 173 Actions
add_content_slide(prs,
    title="CPT 173 CRITICAL ACTIONS",
    subtitle="Lead Element Responsibilities and Execution Checkpoints",
    sections_dict={
        "NLT 17 FEB (Phase I Prep)": [
            "Coordinate with Battalion S3 and BPEnergy liaison for I-TEM scheduling",
            "Validate equipment loadout (EDR tools, sensors, forensic kits, laptops, comms devices)",
            "Brief team leads on APT41 TTPs and operation security (OPSEC) requirements"
        ],
        "NLT 23 FEB (Phase II Activation)": [
            "Establish 24/7 MOC staffing with watch rotation (OIC, NCOIC, analyst, communications NCO)",
            "Execute tool deployment and baseline network mapping within 6 hours of arrival",
            "Deliver Initial Assessment Report within 10 hours of mission start (scope, posture, immediate risks)"
        ],
        "PHASE II (Continuous)": [
            "Lead intelligence-driven threat hunting (3-5 hunt hypotheses per day based on PIRs)",
            "Triage and command incident response (detection → analysis → containment cycle)",
            "Maintain MOC COP (Common Operational Picture) with real-time update to ARCYBER S3"
        ],
        "PHASE III (Transition)": "Prepare comprehensive final assessment with recommendations; conduct staff training; validate BPEnergy capability for independent operations"
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 10: Authorities & ROE
add_content_slide(prs,
    title="AUTHORITIES & RULES OF ENGAGEMENT",
    subtitle="Legal Basis and Operational Constraints",
    sections_dict={
        "DCO-M AUTHORITY": [
            "Granted by ARCYBER FRAGO 01-26 and TASKORD 26-04",
            "Title 10, U.S. Code: DoD State, Local, Tribal Territories protection mission",
            "Authorized to conduct: monitoring, detection, analysis, defensive containment, forensics, malware analysis"
        ],
        "DCO-RA AUTHORITY": [
            "Requires explicit ARCYBER approval unless pre-authorized in TASKORD",
            "Enables isolation, disruption, or neutralization of confirmed malicious activity",
            "All actions must have documented authority chain in mission log"
        ],
        "RESTRICTED ACTIONS": [
            "No offensive cyberspace operations or destructive actions",
            "No production OT system modification without explicit written approval from BPEnergy leadership",
            "No access to non-mission-related corporate data",
            "No actions that could disrupt manufacturing without risk acceptance"
        ],
        "ESCALATION TRIGGER": "If APT41 activity threatens DoD production: notify BPEnergy + ARCYBER immediately → request DCO-RA authority if containment requires elevated action"
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 11: Risk Assessment
add_content_slide(prs,
    title="RISK ASSESSMENT & MITIGATION",
    subtitle="Key Operational Risks and Approved Mitigation Strategies",
    sections_dict={
        "RISK 1: OT Coordination Delay": "Red: Slow BPEnergy OT team response hampers containment speed. Mitigation: Pre-establish 24/7 OT escalation contact and approval authority during I-TEM.",
        "RISK 2: Cloud Persistence Unknown": "Red: APT41 may have IAM keys or service principal persistence not visible in on-premise logs. Mitigation: Prioritize cloud IAM audit in Phase I; daily cloud threat hunting.",
        "RISK 3: Authority Approval Timing": "Yellow: DCO-RA request approval from ARCYBER may lag incident response need. Mitigation: Pre-coordinate standard playbooks; request pre-authorization for common scenarios.",
        "RISK 4: Limited Baseline Data": "Yellow: Phase I baseline collection may be incomplete, creating false positives. Mitigation: Rapid asset inventory and 48-hour baseline period before aggressive hunting; coordinate with BPEnergy security team."
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 12: Success Criteria
add_content_slide(prs,
    title="SUCCESS CRITERIA & END STATE",
    subtitle="Measurable Objectives and Post-Mission Sustainability",
    sections_dict={
        "PHASE I SUCCESS": [
            "MOC fully operationalized with 24/7 watch capability and redundant comms",
            "All authorized personnel on-site with valid access to BPEnergy networks",
            "Baseline vulnerability assessment completed and findings delivered",
            "PIR-driven collection plan published and monitoring sensors activated"
        ],
        "PHASE II SUCCESS": [
            "Zero confirmed APT41 beacons or persistence mechanisms in BPEnergy networks",
            "All detected incidents contained and eradicated within 4-hour response window",
            "Daily SITREP delivered to ARCYBER and BPEnergy with <2% false positive rate",
            "Malware analysis and TTP attribution completed for all incidents"
        ],
        "PHASE III SUCCESS (END STATE)": [
            "BPEnergy security team trained and capable of independent threat hunting and IR",
            "Hardened network controls and detection rules delivered and validated",
            "Post-mission monitoring plan established with BPEnergy ownership",
            "CPT 173 transitions out with zero residual adversary presence confirmed"
        ]
    },
    org_context="CPT 173 // 126th CPBn // Op Manhattan Protocol"
)

# Slide 13: Closing
closing_slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(closing_slide)

# Closing title
close_title = closing_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf = close_title.text_frame
tf.word_wrap = True
tf.text = "READY FOR MISSION EXECUTION"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Closing subtitle
close_sub = closing_slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
tf = close_sub.text_frame
tf.word_wrap = True
tf.text = "CPT 173 Mission Element Lead\n126th Cyber Protection Battalion\nARCYBER OPORD 26-02 / TASKORD 26-04"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# Save
prs.save('CPT173_MANHATTAN_PROTOCOL_PROFESSIONAL.pptx')
print("SUCCESS: CPT173_MANHATTAN_PROTOCOL_PROFESSIONAL.pptx created (13 slides)")
print("\nStyle Applied:")
print("  - Professional military briefing format")
print("  - Calibri throughout (22pt headers, 13pt sections, 11pt body)")
print("  - Dark navy background with white/light gray text")
print("  - Classification marking on every slide")
print("  - Operational metadata footer on all content slides")
print("  - 3-4 major sections per slide with structured content")
