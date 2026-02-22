from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 40, 80)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_content_slide(prs, title, content_dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)

    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0, 40, 80)
    title_shape.line.color.rgb = RGBColor(0, 40, 80)

    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.vertical_anchor = 1

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.8)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for key, value in content_dict.items():
        p = text_frame.add_paragraph()
        p.text = key
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 40, 80)
        p.space_before = Pt(6)
        p.space_after = Pt(3)

        if isinstance(value, list):
            for item in value:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(13)
                p.level = 1
                p.space_after = Pt(2)
        else:
            p = text_frame.add_paragraph()
            p.text = value
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(2)

    return slide

# Slide 1: Cover
add_title_slide(prs, "CPT173 MDMP BRIEFING",
                "OPORD 26-02: Defensive Cyberspace Operations\nBPEnergy Critical Infrastructure Defense")

# Slide 2: MDMP Overview
add_content_slide(prs, "MDMP — Seven-Step Process", {
    "Mission Context": "Deploy NLT 23 FEB 2026 to conduct DCO-M and DCO-RA at BPEnergy in support of DoD state-local-tribal territories protection.",
    "MDMP Steps": [
        "1. Receipt of Mission",
        "2. Mission Analysis",
        "3. COA Development",
        "4. COA Analysis (Wargame)",
        "5. COA Comparison",
        "6. COA Approval",
        "7. Orders Production"
    ]
})

# Slide 3: Step 1
add_content_slide(prs, "STEP 1: RECEIPT OF MISSION", {
    "Higher HQ Direction": [
        "ARCYBER FRAGO 01-26 to OPORD 25-01",
        "USCYBERCOM TASKORD 26-04"
    ],
    "Mission Type": "Defensive Cyberspace Operations (DCO-M / DCO-RA)",
    "Supported Organization": "Town of BPEnergy Cybersecurity Division (125K citizens, CI facilities)",
    "Timeline": "Deployment NLT 23 FEB 2026; Duration TBD",
    "Initial RFIs": [
        "Network architecture and access points",
        "Current EDR/SIEM posture",
        "Current adversary indicators/TTPs"
    ]
})

# Slide 4: Step 2A - OE
add_content_slide(prs, "STEP 2: MISSION ANALYSIS — Operational Environment", {
    "Area of Operations": [
        "Primary: BPEnergy Enterprise Network (SI-EN)",
        "Secondary: Bakerville Cloud Environments (Azure Gov, AWS GovCloud)",
        "Tertiary: OT enclaves supporting critical infrastructure"
    ],
    "Threat": "APT41 (state-sponsored, PRC-attributed)",
    "MLCOA": "Credential harvesting, lateral movement, data exfiltration targeting CI OT",
    "MDCOA": "Destructive malware targeting OT systems to disrupt public services"
})

# Slide 5: Step 2B - Task/Purpose
add_content_slide(prs, "STEP 2: MISSION ANALYSIS — Essential Tasks", {
    "Task": "Conduct DCO-M and DCO-RA; lead enterprise-wide vulnerability assessment and threat hunting",
    "Purpose": "Identify, contain, and mitigate malicious cyber activity targeting BPEnergy critical networks",
    "End State": [
        "BPEnergy networks hardened",
        "Adversary presence removed or contained",
        "Critical services sustain secure operations"
    ],
    "Success Criteria": [
        "Initial Assessment within 10 hours",
        "Daily SITREPs at 1600Z",
        "Incident Reports within 1 hour"
    ]
})

# Slide 6: Step 2C - PIRs
add_content_slide(prs, "STEP 2: MISSION ANALYSIS — Priority Intelligence Requirements", {
    "PIR 1": "Has APT41 established persistence within BPEnergy enterprise or cloud?",
    "PIR 2": "What TTPs is APT41 currently employing against enterprise and CI targets?",
    "PIR 3": "Is APT41 attempting to access or disrupt OT systems supporting CI?",
    "PIR 4": "What indicators reveal imminent lateral movement or exfiltration?",
    "Collection": "SIEM queries, EDR telemetry, network baselines, threat intel fusion"
})

# Slide 7: Step 3A - Task Org
add_content_slide(prs, "STEP 3: COA DEVELOPMENT — Task Organization", {
    "CPT 173 (Lead Element)": [
        "Establish and operate MOC",
        "Enterprise vulnerability assessment",
        "Cloud/hybrid threat hunting",
        "Lead incident response and DCO-RA execution"
    ],
    "CPT 174 (Support)": [
        "Malware analysis and reverse engineering",
        "Surge support for DCO-RA actions"
    ],
    "Battalion S2": "Daily INTSUMs, threat actor profiles",
    "Supporting": "Battalion S3, 91st CPB (reachback), DHS CISA, FBI Cyber"
})

# Slide 8: Step 3B - Concept
add_content_slide(prs, "STEP 3: COA DEVELOPMENT — Concept of Operations", {
    "Phase I (17-22 FEB)": [
        "Conduct I-TEM with BPEnergy",
        "Establish MOC and secure comms",
        "Validate ROE, authorities, access"
    ],
    "Phase II (23 FEB-TBD)": [
        "Continuous monitoring and threat hunting",
        "DCO-RA execution on confirmed activity",
        "Malware analysis, incident response"
    ],
    "Phase III (TBD)": [
        "Final reports and mitigation plans",
        "Knowledge transfer to BPEnergy staff"
    ]
})

# Slide 9: Step 3C - Synch
add_content_slide(prs, "STEP 3: COA DEVELOPMENT — Synchronization", {
    "MOC Operation": "24/7 command and control, COP maintenance",
    "PACE Plan": [
        "Primary: NIPR via ARCYBER-approved VPN",
        "Alternate: BPEnergy secure collaboration environment",
        "Contingency: Encrypted mobile devices",
        "Emergency: Voice-only through BPEnergy security office"
    ],
    "Reporting": [
        "Daily SITREP: 1600Z",
        "Incident reports: Within 1 hour",
        "Weekly summary: Thursday 1600Z"
    ]
})

# Slide 10: Step 4A - Wargame
add_content_slide(prs, "STEP 4: COA ANALYSIS — Wargame Framework", {
    "Scenario 1: Credential Compromise": [
        "Red COA: Spearphish with MFA bypass",
        "Blue Response: Disable accounts, reset credentials",
        "Decision Points: Scope of lateral movement?"
    ],
    "Scenario 2: Malware Infection": [
        "Red COA: ShadowPad or Winnti beacon deployment",
        "Blue Response: Isolate host, analyze, eradicate",
        "Decision Points: C2 communication patterns?"
    ],
    "Scenario 3: OT Attack": [
        "Red COA: Destructive malware on OT network",
        "Blue Response: Coordinate with BPEnergy OT team",
        "Decision Points: Risk to manufacturing?"
    ]
})

# Slide 11: Step 4B - Friction
add_content_slide(prs, "STEP 4: COA ANALYSIS — Friction Points & Shortfalls", {
    "Friction 1 – OT Coordination": "Slow response if OT team unavailable; Mitigation: Pre-establish escalation contact.",
    "Friction 2 – Cloud Visibility": "Unknown persistence in cloud; Mitigation: Prioritize cloud IAM audit early.",
    "Friction 3 – Authority Delays": "DCO-RA approval may delay containment; Mitigation: Pre-coordinate playbooks with ARCYBER.",
    "Shortfall": "Limited baseline data at mission start; Mitigation: Rapid asset inventory in Phase I.",
    "Risk Acceptance": "No production changes without approval; owner: BPEnergy leadership"
})

# Slide 12: Step 5 - Comparison
add_content_slide(prs, "STEP 5: COA COMPARISON — Evaluation Criteria", {
    "Feasibility": "Do we have tools, access, and authority? YES (all pre-arranged via I-TEM).",
    "Acceptability": "Risk to BPEnergy? MITIGATION: ROE limits production changes; must pre-coordinate.",
    "Sustainability": "Can BPEnergy sustain post-mission? YES, Phase III includes knowledge transfer.",
    "Time": "Can we respond fast enough? YES, 1-hour incident reporting pre-negotiated.",
    "Recommended COA": "Integrated DCO-M/DCO-RA with distributed threat hunting and centralized MOC."
})

# Slide 13: Step 6 - Approval
add_content_slide(prs, "STEP 6: COA APPROVAL — Commander's Decision", {
    "Battalion Commander": "LTC Jackson approves three-phase operational concept.",
    "Mission OIC": "MAJ Manbear validates task organization and reporting requirements.",
    "ARCYBER Approval": "FRAGO 01-26 and TASKORD 26-04 authorize DCO-M and pre-authorized DCO-RA.",
    "BPEnergy Consent": "Access Agreement signed; ROE established; escalation procedures confirmed.",
    "Approval Authority": "Title 10 authorities; DoD State-Local-Tribal Territories protection mandate."
})

# Slide 14: Step 7 - Orders
add_content_slide(prs, "STEP 7: ORDERS PRODUCTION — Cyber Annex Products", {
    "OPORD 26-02 Annexes": [
        "Annex A (Task Organization)",
        "Annex B (Intelligence)",
        "Annex C (Operations)"
    ],
    "Supporting Annexes": [
        "Annex J (Cyber Technical Procedures)",
        "Annex K (Incident Response Playbooks)",
        "Annex L (Rules of Engagement / Authorities)"
    ],
    "Subordinate Tasks": [
        "CPT 173: MOC establishment, enterprise assessment, lead IR",
        "CPT 174: Malware analysis, surge support",
        "S2: Daily INTSUMs, threat actor profiles",
        "S3: Synchronization, logistics, reporting"
    ]
})

# Slide 15: Execution
add_content_slide(prs, "EXECUTION FRAMEWORK — CPT173 Responsibilities", {
    "Phase I": [
        "I-TEM coordination (17-22 FEB)",
        "MOC setup and staffing",
        "Network baseline and asset inventory"
    ],
    "Phase II": [
        "24/7 monitoring and threat hunting",
        "Incident response and DCO-RA execution",
        "Daily SITREP and threat updates"
    ],
    "Phase III": [
        "Final assessment and recommendations",
        "Staff training and knowledge transfer",
        "Validation of security hardening"
    ]
})

# Slide 16: Timeline
add_content_slide(prs, "KEY DATES & TIMELINE", {
    "17 FEB 2026": "Phase I begins – I-TEM and MOC establishment",
    "23 FEB 2026": "Phase II begins – Active monitoring (NLT)",
    "10 Hours": "Initial Assessment Report deadline from mission start",
    "1600Z Daily": "Daily SITREP to ARCYBER and BPEnergy leadership",
    "1 Hour": "Incident Report deadline from detection",
    "TBD": "Phase III transition and hardening (commander's decision)"
})

# Slide 17: Closing
add_title_slide(prs, "QUESTIONS", "CPT173 Ready for Mission Execution\nARCYBER OPORD 26-02 / TASKORD 26-04")

prs.save('CPT173_MDMP_PHASES_BRIEF.pptx')
print("CPT173_MDMP_PHASES_BRIEF.pptx created successfully")
